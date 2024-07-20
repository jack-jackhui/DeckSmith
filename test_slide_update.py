from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from slide_deck_generator import remove_unwanted_slides

# Path to your PowerPoint template
template_path = 'slide_template/template-1.pptx'
output_path = 'test_output_presentation.pptx'

# Load the presentation
prs = Presentation(template_path)

# Indices of slides to remove (example: removing slides 0, 2, and 4)
indices_to_remove = list(range(10, len(prs.slides)))

# Function to modify the content of a specific slide
def modify_slide(prs, slide_index, title_text, body_text):
    # Select the slide based on the index
    slide = prs.slides[slide_index]

    # Modify the title
    title_placeholder = slide.shapes.title
    if title_placeholder:
        title_placeholder.text = title_text

    # Modify the body text
    body_placeholder = slide.placeholders[3]  # Assuming the body text is in placeholder 1
    if body_placeholder:
        text_frame = body_placeholder.text_frame
        text_frame.clear()  # Clear existing content

        # Add new paragraphs
        p = text_frame.add_paragraph()
        p.text = body_text
        p.font.size = Pt(16)  # Adjust the font size as needed
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT  # Align text to the left


# Modify the content of the first slide (index 0)
modify_slide(prs, slide_index=5, title_text="New Title", body_text="This is the new body text.")
# Remove the specified slides
prs = remove_unwanted_slides(prs, indices_to_remove)
# Save the modified presentation
prs.save(output_path)

print(f"Presentation saved to {output_path}")