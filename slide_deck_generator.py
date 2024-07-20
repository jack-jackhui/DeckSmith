from langchain_openai import AzureChatOpenAI
import os
from dotenv import load_dotenv
from pptx import Presentation, enum
from pptx.util import Pt, Inches, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from PIL import Image
from image_generator import generate_cover

# Load environment variables from .env file
load_dotenv()

# Initialize the Azure OpenAI client
azure_openai_api_key = os.getenv("AZURE_OPENAI_API_KEY")
azure_openai_endpoint = os.getenv("AZURE_OPENAI_ENDPOINT")
azure_openai_api_version = os.getenv("AZURE_OPENAI_API_VERSION")
azure_openai_deployment_name = os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME")
azure_openai_model_name = os.getenv("AZURE_OPENAI_MODEL_NAME")

llm = AzureChatOpenAI(
    api_key=azure_openai_api_key,
    azure_endpoint=azure_openai_endpoint,
    api_version=azure_openai_api_version,
    deployment_name=azure_openai_deployment_name,
    model_name=azure_openai_model_name
)

# Define the structure of your presentation with layouts and titles
slides_structure = [
    {"layout": "title", "title": "Front Page"},
    {"layout": "title_and_body", "title": "Executive Summary"},
    {"layout": "title_and_body", "title": "Key Point 1"},
    {"layout": "title_and_body", "title": "Key Point 2"},
    {"layout": "title_and_body", "title": "Key Point 3"},
    {"layout": "title_and_body", "title": "Recommendations"},
    {"layout": "title_and_body", "title": "Conclusion"},
    {"layout": "title_and_body", "title": "Additional Examples"},
    {"layout": "title_and_body", "title": "Important Data"},
    {"layout": "title_and_body", "title": "Benefits"},
    {"layout": "title_and_body", "title": "Risks"},
    {"layout": "title_and_body", "title": "Final Conclusion"}
]

layout_mapping = {
    "title": 0,
    "section_header": 1,
    "title_and_body": 2,
    "title_and_two_columns": 3,
    "title_only": 4,
    "one_column_text": 5,
    "main_point": 6,
    "section_title_and_description": 7,
    "caption_only": 8,
    "big_number": 9,
    "blank": 10
}


def remove_unwanted_slides(prs, indices_to_remove):
    # Ensure the indices are sorted in descending order
    indices_to_remove = sorted(indices_to_remove, reverse=True)

    for index in indices_to_remove:
        slide_id = prs.slides._sldIdLst[index].rId
        prs.part.drop_rel(slide_id)
        del prs.slides._sldIdLst[index]

    return prs


def clear_existing_slides(prs, last_index_to_clear):
    while len(prs.slides) > 0:
        # Remove the first slide until all are gone
        # This ensures all associated files and relationships are also cleared
        xml_slides = prs.slides._sldIdLst
        prs.part.drop_rel(xml_slides[0].rId)  # Drop relationship to clean up properly
        del xml_slides[0]


def generate_slide_content(description):
    # Single prompt for generating the entire slide deck content
    prompt = (
        f"You are an expert in creating professional presentations in the style of McKinsey, BCG, or Bain (MBB). "
        f"Generate concise, impactful content for a PowerPoint presentation based on the following description:\n\n"
        f"{description}\n\n"
        "Structure the content for the following slides. Each slide's content should be separated by the marker '[SLIDE]'.\n"
        "1. Front Page: Include a concise title (the description), a sub-headline, the company's name, and the date.\n"
        "2. Executive Summary: Provide a brief summary using the Situation-Complication-Resolution framework, presenting the initial context (Situation), the specific challenge (Complication), and the proposed solution (Resolution).\n"
        "3. Key Point 1: This slide presents key point 1 contains two main parts:an concise Action Title that articulates the key takeaway or main message whic capture the audience's attention, explain the page's importance, and show its contribution to the storyline; Slide Body that allocated to content and exhibits that support the action title or headlines. These exhibits can take the form of graphs, tables, text, maps, or other forms of data visualization. The focus should be on presenting value-adding insights and synthesizing information.\n"
        "4. Key Point 2: This slide presents key point 2 contains two main parts:an concise Action Title that articulates the key takeaway or main message whic capture the audience's attention, explain the page's importance, and show its contribution to the storyline; Slide Body that allocated to content and exhibits that support the action title or headlines. These exhibits can take the form of graphs, tables, text, maps, or other forms of data visualization. The focus should be on presenting value-adding insights and synthesizing information.\n"
        "5. Key Point 3: This slide presents key point 3 contains two main parts:an concise Action Title that articulates the key takeaway or main message whic capture the audience's attention, explain the page's importance, and show its contribution to the storyline; Slide Body that allocated to content and exhibits that support the action title or headlines. These exhibits can take the form of graphs, tables, text, maps, or other forms of data visualization. The focus should be on presenting value-adding insights and synthesizing information.\n"
        "6. Recommendations: Provide concise recommendations outlines the actions required to address the problem or challenge discussed earlier. To create effective recommendations, it is beneficial to group them into categories for clarity, label or number the groups and individual recommendations, and use an active voice starting with action words (verbs) to make the recommendations actionable and compelling.\n"
        "7. Conclusion: Summarize key points, recommendations, and actions required concisely.\n"
        "8. Additional Examples: Provide concise additional examples or detailed analysis with action titles and subheadings.\n"
        "9. Important Data: Highlight important data or statistics concisely with action titles and subheadings.\n"
        "10. Benefits: Emphasize benefits concisely with action titles and subheadings.\n"
        "11. Risks: Outline potential risks and mitigation strategies concisely with action titles and subheadings.\n"
        "12. Final Conclusion: Provide a final conclusion and call to action concisely with action titles and subheadings.\n"
        "Ensure the content is engaging, informative, and suitable for a professional audience. "
        "Include relevant details and examples. Keep the content concise and within the size of the slide. "
        "Do not write any comments other than actual contents of the slide."
    )

    response = llm.invoke(prompt)
    content = response.content.strip()
    print(f"Generated Content:\n{content}\n")

    # Split the content into separate slides based on the marker
    slides_content = content.split("[SLIDE]")
    for index, slide in enumerate(slides_content):
        print(f"Content for Slide {index + 1}:\n{slide}\n")
    slides_content = [slide.strip() for slide in slides_content if slide.strip()]

    # Debug: Print the number of slides content generated
    print(f"Number of slides content generated: {len(slides_content)}")
    print(f"Expected number of slides: {len(slides_structure)}")

    # Ensure slides_content aligns with slides_structure
    if len(slides_content) != len(slides_structure):
        raise ValueError("Generated content does not match the number of defined slides")

    return slides_content


def adjust_slide_content(slide, title, content, resize_title=True, resize_content=True):
    # Adjust title
    if resize_title:
        title_placeholder = slide.shapes.title
        if title_placeholder:
            title_placeholder.top = Inches(1.5)  # Move title to the top
            title_placeholder.width = Inches(10)  # Set width
            title_placeholder.height = Inches(1.5)  # Set height

    # Adjust content
    if resize_content:
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Typically the content placeholder
                content_placeholder = shape
                break
        if content_placeholder:
            content_placeholder.top = Inches(1.5)  # Move content box down
            content_placeholder.left = Inches(0.5)  # Adjust left margin
            content_placeholder.width = Inches(9)  # Adjust width
            content_placeholder.height = Inches(5.5)  # Adjust height

    insert_content(slide, title, content)


def create_slide(prs, layout_name):
    layout_index = layout_mapping.get(layout_name)
    if layout_index is None:
        raise ValueError(f"No layout found for {layout_name}")
    slide_layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(slide_layout)


"""
def insert_content(slide, title, content):
    # Set title
    try:
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
    except AttributeError:
        print("Warning: No title placeholder found. Skipping title.")

    # Find the content placeholder and set content
    try:
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Typically the content placeholder
                content_placeholder = shape
                break
        if content_placeholder:
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # Clear existing content
            p = text_frame.add_paragraph()
            p.text = content
            p.font.size = Pt(16)  # Adjust the font size as needed
            # Adjust position if this is a slide with an image
            if slide.slide_layout == layout_mapping["one_column_text"]:
                content_placeholder.top = Inches(0.5)  # Move text box up
                content_placeholder.width = Inches(4.5)  # Adjust width for left column
        else:
            print("Warning: No content placeholder found. Skipping content insertion.")
    except Exception as e:
        print(f"Error while inserting content: {str(e)}")
"""


def insert_content(slide, title, content, image_slide=False):
    try:
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
    except AttributeError:
        print("Warning: No title placeholder found. Skipping title.")

    try:
        # For one_column_text layout, use the correct text placeholder index
        content_placeholder = slide.placeholders[1]

        text_frame = content_placeholder.text_frame
        text_frame.clear()  # Clear any existing content
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(16)  # Adjust font size as necessary
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        # Adjust text box width if it's a slide with an image
        if image_slide:
            content_placeholder.width = Inches(4.5)  # Reduce width to avoid overlap with image
    except Exception as e:
        print(f"Error setting text: {str(e)}")


def add_image_to_slide(slide, image_path):
    try:
        image = Image.open(image_path)
        image_width, image_height = image.size

        image_width_cm = image_width * 0.0264583333
        image_height_cm = image_height * 0.0264583333

        slide_width_cm = 24.4
        slide_height_cm = 19.05

        # Resize the image to fit well on the right side, taking up no more than one-third of the slide
        while image_width_cm > slide_width_cm / 3 or image_height_cm > slide_height_cm:
            image_width_cm *= 0.5
            image_height_cm *= 0.5

        top = Cm(2.0)
        left = Cm(slide_width_cm - image_width_cm - 1.0)  # Adjust to position the image to the right side
        height = Cm(image_height_cm)
        width = Cm(image_width_cm)

        slide.shapes.add_picture(image_path, left=left, top=top, width=width, height=height)
    except Exception as e:
        print(f"Error adding image to slide: {str(e)}")

def generate_and_save_presentation(description, template_path, output_path):
    slide_content = generate_slide_content(description)
    prs = Presentation(template_path)
    clear_existing_slides(prs, 53)

    slides_with_images = [1, 6, 7]  # Executive Summary, Conclusion, Additional Examples

    for i, slide_data in enumerate(slides_structure):
        slide = create_slide(prs, slide_data["layout"])
        title = description if i == 0 else slide_data["title"]

        if i in slides_with_images:
            insert_content(slide, title, slide_content[i], image_slide=True)
            module = {"topic": title, "key_points": slide_content[i].split('\n')}
            module = generate_cover(i, module)
            if module["filename"]:
                add_image_to_slide(slide, module["filename"])
        else:
            insert_content(slide, title, slide_content[i])

        print(f"Inserted content for slide {i + 1}")

    prs.save(output_path)
    print(f"Presentation saved to {output_path}")
