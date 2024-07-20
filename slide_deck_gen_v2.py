from langchain_openai import AzureChatOpenAI
import os
from dotenv import load_dotenv
from pptx import Presentation, enum
from pptx.util import Pt, Inches, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from PIL import Image
from image_generator import generate_cover

# Load environment variables from .env file
load_dotenv()

# Initialize the Azure OpenAI client
azure_openai_api_key = os.getenv("AZURE_OPENAI_API_KEY")
azure_openai_endpoint = os.getenv("AZURE_OPENAI_ENDPOINT")
azure_openai_api_version = os.getenv("AZURE_OPENAI_API_VERSION")
azure_openai_deployment_name = os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME")
azure_openai_model_name = os.getenv("AZURE_OPENAI_MODEL_NAME")

llm = AzureChatOpenAI(
    api_key=azure_openai_api_key,
    azure_endpoint=azure_openai_endpoint,
    api_version=azure_openai_api_version,
    deployment_name=azure_openai_deployment_name,
    model_name=azure_openai_model_name
)

slides_structure = [
    {"title": "Front Page"},
    {"title": "Executive Summary"},
    {"title": "Key Point 1"},
    {"title": "Key Point 2"},
    {"title": "Key Point 3"},
    {"title": "Recommendations"},
    {"title": "Conclusion"},
    {"title": "Additional Examples"},
    {"title": "Important Data"},
    {"title": "Benefits"},
    {"title": "Risks"},
    {"title": "Final Conclusion"}
]

layout_mapping = {
    "title": 0,
    "section_header": 1,
    "title_and_body": 2,
    "title_and_two_columns": 3,
    "title_only": 4,
    "one_column_text": 5,
    "main_point": 6,
    "section_title_and_description": 7,
    "caption_only": 8,
    "big_number": 9,
    "blank": 10
}

def remove_unwanted_slides(prs, indices_to_remove):
    # Ensure the indices are sorted in descending order
    indices_to_remove = sorted(indices_to_remove, reverse=True)

    for index in indices_to_remove:
        slide_id = prs.slides._sldIdLst[index].rId
        prs.part.drop_rel(slide_id)
        del prs.slides._sldIdLst[index]

    return prs


def clear_existing_slides(prs, last_index_to_clear):
    while len(prs.slides) > 0:
        # Remove the first slide until all are gone
        # This ensures all associated files and relationships are also cleared
        xml_slides = prs.slides._sldIdLst
        prs.part.drop_rel(xml_slides[0].rId)  # Drop relationship to clean up properly
        del xml_slides[0]

def generate_slide_content(description):
    # Single prompt for generating the entire slide deck content
    prompt = (
        f"You are an expert in creating professional presentations in the style of McKinsey, BCG, or Bain (MBB). "
        f"Generate concise, impactful content for a PowerPoint presentation based on the following description:\n\n"
        f"{description}\n\n"
        "Structure the content for the following slides. Each slide's content should be separated by the marker '[SLIDE]'.\n"
        "Do not use any markdown or special formatting syntax in the content.\n"
        "Do not repeat the slide titles in the body text for any slides.\n"
        "1. Front Page: Include a concise title (the description), and a catchy sub-headline.\n"
        "2. Executive Summary: Provide a brief summary using the Situation-Complication-Resolution framework, presenting the initial context (Situation), the specific challenge (Complication), and the proposed solution (Resolution).\n"
        "3. Key Point 1: This slide presents key point 1 contains two main parts:an concise Action Title that articulates the key takeaway or main message whic capture the audience's attention, explain the page's importance, and show its contribution to the storyline; Slide Body that allocated to content and exhibits that support the action title or headlines. These exhibits can take the form of graphs, tables, text, maps, or other forms of data visualization. The focus should be on presenting value-adding insights and synthesizing information.\n"
        "4. Key Point 2: This slide presents key point 2 contains two main parts:an concise Action Title that articulates the key takeaway or main message whic capture the audience's attention, explain the page's importance, and show its contribution to the storyline; Slide Body that allocated to content and exhibits that support the action title or headlines. These exhibits can take the form of graphs, tables, text, maps, or other forms of data visualization. The focus should be on presenting value-adding insights and synthesizing information.\n"
        "5. Key Point 3: This slide presents key point 3 contains two main parts:an concise Action Title that articulates the key takeaway or main message whic capture the audience's attention, explain the page's importance, and show its contribution to the storyline; Slide Body that allocated to content and exhibits that support the action title or headlines. These exhibits can take the form of graphs, tables, text, maps, or other forms of data visualization. The focus should be on presenting value-adding insights and synthesizing information.\n"
        "6. Recommendations: Provide concise recommendations in the form of three key points, each with a heading and body text. The headings should be impactful and the body text should provide clear, actionable recommendations. Separate each key point with a double newline. Do not include any exhibits or images. Do not include the word recommendation in each point.\n"
        "7. Conclusion: Summarize key points, recommendations, and actions required concisely.\n"
        "8. Additional Examples: Provide two concise examples, each with a heading and body text. Separate each example with a double newline. Do not include any exhibits or images.\n"
        "9. Important Data: Highlight important data or statistics concisely with action titles and subheadings.\n"
        "10. Benefits: Emphasize benefits concisely with action titles and subheadings.\n"
        "11. Risks: Outline potential risks and mitigation strategies concisely with action titles and subheadings.\n"
        "12. Final Conclusion: Provide a final conclusion and call to action concisely with action titles and subheadings.\n"
        "Ensure the content is engaging, informative, and suitable for a professional audience. "
        "Include relevant details and examples. Keep the content concise and within the size of the slide. "
        "Do not write any comments other than actual contents of the slide."
    )

    response = llm.invoke(prompt)
    content = response.content.strip()
    print(f"Generated Content:\n{content}\n")

    # Split the content into separate slides based on the marker
    slides_content = content.split("[SLIDE]")
    for index, slide in enumerate(slides_content):
        print(f"Content for Slide {index + 1}:\n{slide}\n")
    slides_content = [slide.strip() for slide in slides_content if slide.strip()]

    # Debug: Print the number of slides content generated
    print(f"Number of slides content generated: {len(slides_content)}")
    print(f"Expected number of slides: {len(slides_structure)}")

    # Ensure slides_content aligns with slides_structure
    if len(slides_content) != len(slides_structure):
        raise ValueError("Generated content does not match the number of defined slides")

    return slides_content


def adjust_slide_content(slide, title, content, resize_title=True, resize_content=True):
    # Adjust title
    if resize_title:
        title_placeholder = slide.shapes.title
        if title_placeholder:
            title_placeholder.top = Inches(1.5)  # Move title to the top
            title_placeholder.width = Inches(10)  # Set width
            title_placeholder.height = Inches(1.5)  # Set height

    # Adjust content
    if resize_content:
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Typically the content placeholder
                content_placeholder = shape
                break
        if content_placeholder:
            content_placeholder.top = Inches(1.5)  # Move content box down
            content_placeholder.left = Inches(0.5)  # Adjust left margin
            content_placeholder.width = Inches(9)  # Adjust width
            content_placeholder.height = Inches(5.5)  # Adjust height

    insert_content(slide, title, content)

def create_slide(prs, layout_name):
    layout_index = layout_mapping.get(layout_name)
    if layout_index is None:
        raise ValueError(f"No layout found for {layout_name}")
    slide_layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(slide_layout)

"""
def insert_content(slide, title, content):
    # Set title
    try:
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
    except AttributeError:
        print("Warning: No title placeholder found. Skipping title.")

    # Find the content placeholder and set content
    try:
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Typically the content placeholder
                content_placeholder = shape
                break
        if content_placeholder:
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # Clear existing content
            p = text_frame.add_paragraph()
            p.text = content
            p.font.size = Pt(16)  # Adjust the font size as needed
            # Adjust position if this is a slide with an image
            if slide.slide_layout == layout_mapping["one_column_text"]:
                content_placeholder.top = Inches(0.5)  # Move text box up
                content_placeholder.width = Inches(4.5)  # Adjust width for left column
        else:
            print("Warning: No content placeholder found. Skipping content insertion.")
    except Exception as e:
        print(f"Error while inserting content: {str(e)}")
"""

def insert_content(slide, title, content, text_placeholder_idx, font_size=10, enlarge_text_box=False, image_slide=False, move_title_up=False):
    try:
        title_placeholder = slide.shapes.title
        if title_placeholder:
            title_placeholder.text = title
            if move_title_up:
                title_placeholder.top = Inches(0.5)  # Move title up
    except AttributeError:
        print("Warning: No title placeholder found. Skipping title.")

    try:
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == text_placeholder_idx:
                content_placeholder = shape
                break

        if content_placeholder:
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # Clear any existing content

            # Add new paragraphs
            for paragraph in content.split('\n'):
                p = text_frame.add_paragraph()
                p.text = paragraph
                p.font.size = Pt(font_size)  # Adjust font size as necessary

            # Adjust text box width if it's a slide with an image
            if image_slide:
                content_placeholder.width = Inches(4.5)  # Reduce width to avoid overlap with image
        else:
            print(f"Warning: No content placeholder found for text insertion with index {text_placeholder_idx}.")
    except Exception as e:
        print(f"Error setting text: {str(e)}")

def add_image_to_slide(slide, image_path):
    try:
        image = Image.open(image_path)
        image_width, image_height = image.size

        image_width_cm = image_width * 0.0264583333
        image_height_cm = image_height * 0.0264583333

        slide_width_cm = 24.4
        slide_height_cm = 19.05

        # Resize the image to fit well on the right side, taking up no more than one-third of the slide
        while image_width_cm > slide_width_cm / 3 or image_height_cm > slide_height_cm:
            image_width_cm *= 0.5
            image_height_cm *= 0.5

        top = Cm(2.0)
        left = Cm(slide_width_cm - image_width_cm - 1.0)  # Adjust to position the image to the right side
        height = Cm(image_height_cm)
        width = Cm(image_width_cm)

        slide.shapes.add_picture(image_path, left=left, top=top, width=width, height=height)
    except Exception as e:
        print(f"Error adding image to slide: {str(e)}")


def regenerate_slide_6_content(description, max_retries=3):
    for attempt in range(max_retries):
        # Regenerate content specifically for slide 6
        prompt = (
            f"You are an expert in creating professional presentations in the style of McKinsey, BCG, or Bain (MBB). "
            f"Generate concise, impactful content for the 'Recommendations' slide based on the following description:\n\n"
            f"{description}\n\n"
            "Provide concise recommendations in the form of three key points, each with a heading and body text. "
            "The headings should be impactful and the body text should provide clear, actionable recommendations. "
            "Separate each key point with a double newline. Do not include any exhibits or images. "
            "Do not include the word recommendation in each point."
        )

        response = llm.invoke(prompt)
        content = response.content.strip()

        key_points = content.split('\n\n')
        if len(key_points) == 3:
            return content
        else:
            print(f"Regeneration attempt {attempt + 1} failed. Retrying...")

    raise ValueError("Regenerated content for slide 6 does not contain exactly 3 key points after multiple attempts")
def generate_and_save_presentation(description, template_path, output_path):
    try:
        slide_content = generate_slide_content(description)
        # Check the content of slide 6
        key_points = slide_content[5].split('\n\n')
        if len(key_points) != 3:
            raise ValueError("Generated content for slide 6 does not contain exactly 3 key points")
    except ValueError as e:
        if "slide 6 does not contain exactly 3 key points" in str(e):
            print("Error in generating content for slide 6. Regenerating...")
            try:
                slide_6_content = regenerate_slide_6_content(description)
                slide_content[5] = slide_6_content
            except ValueError as regen_error:
                print(str(regen_error))
                return
        else:
            raise e
    prs = Presentation(template_path)

    # Ensure there are enough slides in the template
    if len(prs.slides) < len(slides_structure):
        raise ValueError("The template does not contain enough slides to match the structure")

    slides_with_images = []  # Add indices of slides that should have images here
    slides_to_move_title_up = []  # Indices of slides where title should be moved up

    template_type = "template-2" if "template-2" in template_path else "template-1"

    for i, slide in enumerate(prs.slides):
        if i < len(slides_structure):
            title = description if i == 0 else slides_structure[i]["title"]
            content = slide_content[i]

            if i == 5:  # Special case for slide 6 (index 5) with multiple key points
                key_points = content.split('\n\n')  # Assuming key points are separated by double newlines

                if template_type == "template-2":
                    body_indices = [1, 2, 3]
                    heading_indices = [4, 5, 6]
                else:
                    heading_indices = [4, 6, 8]
                    body_indices = [3, 5, 7]

                for j, key_point in enumerate(key_points):
                    if '\n' in key_point:
                        heading, body = key_point.split('\n', 1)
                    else:
                        heading = key_point
                        body = ""
                    insert_content(slide, title, heading, text_placeholder_idx=heading_indices[j], font_size=10,
                                   enlarge_text_box=False)
                    insert_content(slide, title, body, text_placeholder_idx=body_indices[j], font_size=9,
                                   enlarge_text_box=True)
            elif i == 7:  # Special case for slide 8 (index 7) with two examples
                examples = content.split('\n\n')  # Assuming examples are separated by double newlines
                if len(examples) != 2:
                    raise ValueError("Generated content for slide 8 does not contain exactly 2 examples")
                if template_type == "template-2":
                    heading_indices = [3, 4]
                    body_indices = [1, 2]
                else:
                    heading_indices = [2, 4]
                    body_indices = [1, 3]

                for j, example in enumerate(examples):
                    if '\n' in example:
                        heading, body = example.split('\n', 1)
                    else:
                        heading = example
                        body = ""
                    insert_content(slide, title, heading, text_placeholder_idx=heading_indices[j], font_size=10,
                                   enlarge_text_box=False)
                    insert_content(slide, title, body, text_placeholder_idx=body_indices[j], font_size=9,
                                   enlarge_text_box=True)
            elif i in [8, 9, 10]:  # Special case for slides 9, 10, and 11 to replace content
                title_placeholder = slide.shapes.title
                if title_placeholder:
                    title_placeholder.text = slides_structure[i]["title"]
                # Clear existing content
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx != title_placeholder.placeholder_format.idx:
                        shape.text = ""
                # Insert new content
                insert_content(slide, slides_structure[i]["title"], content, text_placeholder_idx=1, font_size=10,
                               enlarge_text_box=True)
            else:
                # Determine the correct text placeholder index
                if template_type == "template-2":
                    text_placeholder_idx = 1
                else:
                    if i in [0, 8, 9, 10]:  # For slides 1, 9, 10, 11
                        text_placeholder_idx = 1
                    else:
                        text_placeholder_idx = 3
                move_title_up = i in slides_to_move_title_up

                if i in slides_with_images:
                    insert_content(slide, title, content, text_placeholder_idx, image_slide=True, move_title_up=move_title_up)
                    module = {"topic": title, "key_points": content.split('\n')}
                    module = generate_cover(i, module)
                    if module["filename"]:
                        add_image_to_slide(slide, module["filename"])
                else:
                    insert_content(slide, title, content, text_placeholder_idx, enlarge_text_box=True, move_title_up=move_title_up)

        print(f"Inserted content for slide {i + 1}")

    # Remove the specified slides
    indices_to_remove = list(range(12, len(prs.slides)))
    prs = remove_unwanted_slides(prs, indices_to_remove)
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")
