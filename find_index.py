from pptx import Presentation

def inspect_placeholders(template_path):
    prs = Presentation(template_path)
    for slide in prs.slides:
        print(f"Inspecting slide {prs.slides.index(slide) + 1}")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.is_placeholder:
                phf = shape.placeholder_format
                print(f"Index: {phf.idx}, Type: {shape.placeholder_format.type}, Text: '{shape.text}'")

# Provide the path to your template
template_path = "slide_template/template-1.pptx"
inspect_placeholders(template_path)